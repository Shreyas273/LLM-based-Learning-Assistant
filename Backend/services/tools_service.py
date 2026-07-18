import os
from fastapi import UploadFile
import PyPDF2
from docx import Document
from pptx import Presentation
from PIL import Image
import io
import asyncio
import json
import re

# NOTE:
# This module previously contained direct Gemini calls and duplicated code.
# It is now a lightweight file-processing helper + wrappers that delegate to
# dedicated services (`summarizer_service`, `solver_service`, `generator_service`).

async def process_file_content(file: UploadFile, file_type: str = "pdf"):
    """Extract text from various file types"""
    content = ""
    try:
        content_bytes = await file.read()
        file_stream = io.BytesIO(content_bytes)

        if file_type == "pdf":
            reader = PyPDF2.PdfReader(file_stream)
            for page in reader.pages:
                text = page.extract_text()
                if text: content += text + "\n"
                
        elif file_type == "docx":
            doc = Document(file_stream)
            for para in doc.paragraphs:
                content += para.text + "\n"
                
        elif file_type == "pptx":
            prs = Presentation(file_stream)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
                        
        elif file_type == "txt":
            content = content_bytes.decode("utf-8")
            
        return content
    except Exception as e:
        print(f"File processing error: {e}")
        return ""

async def summarize_content(content: str, type: str = "general") -> str:
    from services.summarizer_service import summarize_content as _summarize
    res = await _summarize(content, type)
    # summarizer_service returns dict; tools_service historically returned string here
    if isinstance(res, dict):
        return res.get("summary") or ""
    return str(res)


async def analyze_image(file: UploadFile, task: str = "summarize"):
    from services.solver_service import analyze_image as _analyze
    res = await _analyze(file, task)
    if isinstance(res, dict):
        return res.get("result") or res.get("error") or ""
    return str(res)


async def solve_problem(query: str, subject: str) -> str:
    from services.solver_service import solve_problem as _solve
    res = await _solve(query, subject)
    if isinstance(res, dict):
        return res.get("result") or res.get("error") or ""
    return str(res)


async def generate_content(topic: str, task: str):
    from services.generator_service import generate_content as _gen
    res = await _gen(topic, task)
    if isinstance(res, dict):
        return res.get("result") or res.get("error") or ""
    return str(res)
